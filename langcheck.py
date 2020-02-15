"""Module to check if any sentence in a document is Portuguese."""
import os
from os.path import abspath
import win32com.client as win32
import xlrd
import pptx
import csv
import re
from polyglot.detect import Detector
base_path = os.path.dirname(abspath('__file__'))
avail_exts = ['docx', 'doc', 'pptx', 'xls', 'xlsx', 'csv', 'txt', 'rtf']
passed_exts = ['py', 'git', 'spec', 'exe', 'md', 'gitattributes',
               'gitignore', 'zip']


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join(new_output) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(fname, 'r', encoding='utf8')
        txt = text_doc.read()
        text_doc.close()
    elif fname.endswith('.csv'):
        csv_doc = open(fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt


def lbl_langcheck(txt):
    """Check language of document sentence by sentence."""
    tokens = re.split('[?.;!:,\\n\\r]', txt)
    failed_sents = []
    for i in tokens:
        tx = i.strip()
        detector = Detector(tx, quiet=True)
        lan = detector.language.code
        if lan == 'pt':
            failed_sents.append(tx)
    if len(failed_sents) != 0:
        msg = "Some sentences may be in Portuguese. Check following sentences."
    else:
        msg = "No suspected Portuguese sentences found."
    return msg, failed_sents


def final_report(msg, fname, sents=[]):
    """Format for writing to result file."""
    msg_head = '*' * 20 + '\n' + 'Result for {}:'.format(fname) + '\n'
    result_msg = 'RESULT :: ' + msg + '\n'
    if len(sents) != 0:
        sents_msg = "\n**********\n".join(sents)
    else:
        sents_msg = ""
    return msg_head + result_msg + sents_msg


def directory_check(path=base_path):
    """Run whole module on a given directory."""
    msg_list = []
    for i in os.listdir(path):
        if os.path.isdir(i):
            pass
        elif i == 'script_result.txt':
            pass
        elif i.split('.')[-1] in passed_exts:
            pass
        elif i.endswith('.ppt'):
            msg_1 = 'ppt format not supported.\n'
            msg_2 = 'Chen convert {} to pptx and run script again.'.format(i)
            msg_ppt = msg_1 + msg_2
            msg_list.append(final_report(msg_ppt, i))
        elif i.split('.')[-1] in avail_exts:
            trans_text = extract_text(i)
            msg_trans, sentences = lbl_langcheck(trans_text)
            msg_list.append(final_report(msg_trans, i, sentences))
        else:
            null_msg = '{} is not one of the acceptable formats.'.format(i)
            msg_list.append(final_report(null_msg, i))
    result = open('script_result.txt', 'a', encoding='utf8')
    for j in msg_list:
        result.write(j)
    result.close()
    return


if __name__ == '__main__':
    directory_check()
